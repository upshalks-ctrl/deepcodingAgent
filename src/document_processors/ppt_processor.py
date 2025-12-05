"""
PPT处理器
支持PowerPoint文档处理，包括python-pptx和视觉模型两种模式
"""

import logging
from pathlib import Path
from typing import List
from src.models.document import Document
from src.document_processors import create_document_metadata

logger = logging.getLogger(__name__)


class PPTProcessor:
    """PPT文档处理器 - 支持多种处理模式"""

    async def process(self, file_path: str, **kwargs) -> List[Document]:
        """
        处理PPT文档

        Args:
            file_path: PPT文件路径
            **kwargs: 额外参数
                - mode: 处理模式 ('python_pptx', 'vision_memory_buffer', 默认'python_pptx')
                - dpi: 图片DPI（默认200）
                - chunk_after: 是否按页数切分（默认True）
                - combine_pages: 是否合并所有页面（默认True）

        Returns:
            List[Document]: 文档列表
        """
        if not Path(file_path).exists():
            raise FileNotFoundError(f"PPT文件不存在: {file_path}")

        logger.info(f"开始处理PPT: {file_path}")

        # 获取处理模式
        mode = kwargs.get('mode', 'python_pptx')

        # v2模式：内存buffer + base64编码 + 视觉模型
        if mode == 'vision_memory_buffer':
            try:
                return await self._process_with_vision_memory(file_path, **kwargs)
            except ImportError as e:
                logger.warning(f"视觉模型方法依赖缺失: {e}")
            except Exception as e:
                logger.error(f"视觉模型处理失败: {e}")
                # 回退到python-pptx模式
                logger.info("回退到python-pptx模式")
                return await self._process_with_python_pptx(file_path, **kwargs)

        # v1模式：python-pptx
        elif mode == 'python_pptx':
            try:
                return await self._process_with_python_pptx(file_path, **kwargs)
            except ImportError:
                logger.warning("python-pptx未安装")
            except Exception as e:
                logger.error(f"python-pptx处理失败: {e}")

            # 回退到视觉模型
            logger.info("尝试使用视觉模型模式...")
            try:
                return await self._process_with_vision_memory(file_path, **kwargs)
            except Exception as e:
                logger.error(f"视觉模型模式也失败: {e}")

        else:
            raise ValueError(f"不支持的处理模式: {mode}")

        raise RuntimeError("所有PPT处理方法都失败了")

    async def _process_with_python_pptx(self, file_path: str, **kwargs) -> List[Document]:
        """使用python-pptx处理"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装python-pptx: pip install python-pptx")

        logger.info("使用python-pptx处理")

        prs = Presentation(file_path)

        documents = []
        metadata_base = create_document_metadata(
            file_path=file_path,
            document_type='ppt',
            url=f"ppt://{Path(file_path).name}",
            processor='python-pptx',
            total_slides=len(prs.slides)
        )

        for slide_num, slide in enumerate(prs.slides):
            slide_content = []
            slide_title = f"第{slide_num+1}页"

            # 提取幻灯片中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 检查是否是标题（通常第一个大文本框是标题）
                    if len(text) < 100 and slide_num == 0:
                        slide_title = text
                    slide_content.append(text)

            if slide_content:
                content = "\n".join(slide_content)

                # 复制元数据并添加幻灯片特定信息
                metadata = metadata_base.copy()
                metadata['slide_number'] = slide_num + 1
                metadata['slide_title'] = slide_title
                metadata['word_count'] = len(content)

                document = Document(
                    content=content,
                    metadata=metadata,
                    document_type='ppt',
                    source_path=file_path,
                    title=Path(file_path).name,
                    url=f"ppt://{Path(file_path).name}"
                )
                documents.append(document)

        logger.info(f"python-pptx处理完成，共{len(documents)}页")
        return documents

    async def _process_with_vision_memory(self, file_path: str, **kwargs) -> List[Document]:
        """
        使用视觉模型处理PPT v2（python-pptx + 视觉模型）
        将所有幻灯片提取内容，使用视觉模型转换为Markdown

        Args:
            file_path: PPT文件路径
            **kwargs: 额外参数
                - chunk_after: 是否按页数切分（默认True）
                - combine_pages: 是否合并所有页面（默认True）

        Returns:
            List[Document]: 返回按页切分的文档列表
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装python-pptx: pip install python-pptx")

        logger.info("使用视觉模型处理PPT（v2模式 - python-pptx + 视觉模型）")

        # 初始化视觉模型
        llm_vision = None
        try:
            from src.llms.llm import get_llm_by_type
            llm_vision = get_llm_by_type("vision")
            logger.info("已初始化视觉模型")
        except Exception as e:
            logger.warning(f"无法初始化视觉模型: {e}")

        combine_pages = kwargs.get('combine_pages', True)

        # 打开PPT
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        # 创建元数据
        metadata = create_document_metadata(
            processor='python-pptx+Vision',
            total_slides=total_slides
        )

        # 使用python-pptx提取每页内容
        logger.info(f"使用python-pptx提取 {total_slides} 页内容...")

        slides_content = []
        temp_dir = Path("temp_ppt_images")
        temp_dir.mkdir(exist_ok=True)

        for slide_num, slide in enumerate(prs.slides):
            logger.info(f"正在提取第 {slide_num + 1}/{total_slides} 页...")

            slide_info = {
                'slide_number': slide_num + 1,
                'texts': [],
                'tables': [],
                'images': []
            }

            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_info['texts'].append(text)

                # 检查是否有表格
                if shape.shape_type == 19:  # Table
                    try:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        slide_info['tables'].append(table_data)
                    except Exception as e:
                        logger.warning(f"提取表格失败: {e}")

            # 检查是否有图片
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        img_filename = temp_dir / f"slide_{slide_num+1}_img_{len(slide_info['images'])+1}.png"

                        # 保存图片
                        with open(img_filename, 'wb') as f:
                            f.write(image.blob)

                        slide_info['images'].append(str(img_filename))
                    except Exception as e:
                        logger.warning(f"提取图片失败: {e}")

            slides_content.append(slide_info)

        # 如果有视觉模型，使用视觉模型理解内容
        full_markdown = ""
        if llm_vision:
            try:
                logger.info(f"开始使用视觉模型分析 {total_slides} 页...")

                # 构建内容描述
                content_description = []
                for i, slide_info in enumerate(slides_content):
                    slide_desc = [f"### 第{i+1}页"]

                    if slide_info['texts']:
                        slide_desc.append("**文本内容:**")
                        for text in slide_info['texts']:
                            slide_desc.append(f"- {text}")

                    if slide_info['tables']:
                        slide_desc.append("**表格内容:**")
                        for table in slide_info['tables']:
                            # 转换为Markdown表格
                            if table:
                                # 表头
                                headers = " | ".join(table[0])
                                slide_desc.append(f"| {headers} |")
                                # 分隔线
                                separators = " | ".join(["---"] * len(table[0]))
                                slide_desc.append(f"| {separators} |")
                                # 数据行
                                for row in table[1:]:
                                    row_text = " | ".join(row)
                                    slide_desc.append(f"| {row_text} |")

                    if slide_info['images']:
                        slide_desc.append(f"**包含 {len(slide_info['images'])} 张图片**")

                    content_description.append("\n".join(slide_desc))

                # 构建视觉模型输入（文本 + 图片）
                messages = [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": f"""请根据以下PPT内容，生成完整的Markdown文档：

{chr(10).join(content_description)}

要求：
1. 整理和优化文本内容，保持逻辑清晰
2. 将表格转换为标准Markdown表格格式
3. 使用适当的标题层级结构
4. 保持原有的层次和关系
5. 确保输出是纯Markdown格式，不要包含任何解释或前缀
6. 按幻灯片顺序组织内容
7. 在适当位置添加页码标记

请直接返回Markdown格式的文本内容。"""
                            }
                        ]
                    }
                ]

                # 如果有图片，添加到消息中
                for slide_info in slides_content:
                    for img_path in slide_info['images']:
                        messages[0]["content"].append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"file://{Path(img_path).absolute()}"
                            }
                        })  # type: ignore[arg-type]

                response = await llm_vision.ainvoke(messages)

                # 处理响应
                if hasattr(response, 'content'):
                    content_text = response.content
                else:
                    content_text = str(response)

                if isinstance(content_text, list):
                    if content_text and isinstance(content_text[0], dict) and 'content' in content_text[0]:
                        full_markdown = content_text[0]['content']
                    else:
                        full_markdown = str(content_text[0]) if content_text else ""
                else:
                    full_markdown = content_text or ""

                # 如果内容为空，添加占位符
                if not full_markdown.strip():
                    full_markdown = "[PPT内容为空或无法识别]"

                logger.info("视觉模型分析完成")

            except Exception as e:
                logger.error(f"视觉模型分析失败: {e}")
                # 回退到纯文本提取
                full_markdown = await self._extract_text_only(slides_content, total_slides)
        else:
            logger.warning("未配置视觉模型，使用纯文本提取")
            full_markdown = await self._extract_text_only(slides_content, total_slides)

        # 清理临时图片
        for img_file in temp_dir.glob("*.png"):
            img_file.unlink(missing_ok=True)
        if temp_dir.exists() and not any(temp_dir.iterdir()):
            temp_dir.rmdir()

        # 如果需要自动按页切分
        if kwargs.get('chunk_after', True):
            logger.info("开始按页切分Markdown文档...")

            # 直接实现按页切分
            import re
            # 查找页码标记模式：第N页 或 ### 第N页
            page_splits = re.split(r'\n\s*(?:#{0,3}\s*)?第\d+页\s*\n', full_markdown)
            page_splits = [p.strip() for p in page_splits if p.strip()]

            # 如果没有找到明确的页码标记，按段落切分
            if len(page_splits) <= 1:
                paragraphs = re.split(r'\n\s*\n', full_markdown)
                page_splits = [p.strip() for p in paragraphs if p.strip()]

            # 创建按页切分的文档
            chunks = []
            for i, page_content in enumerate(page_splits):
                page_doc = Document(
                    content=page_content,
                    metadata=metadata.copy(),
                    document_type = 'ppt',
                    source_path = file_path,
                    title = f"第{i+1}页",
                    url = f"ppt://{Path(file_path).name}#page={i+1}"
                )
                page_doc.metadata['slide_number'] = i + 1
                page_doc.metadata['chunk_type'] = 'page'
                page_doc.metadata['chunk_title'] = f"第{i+1}页"
                page_doc.metadata['word_count'] = len(page_content)
                page_doc.metadata['processing_method'] = 'python_pptx_vision_combined'

                chunks.append(page_doc)

            logger.info(f"文档切分为 {len(chunks)} 个块（按页）")
            return chunks

        # 创建完整的Markdown文档
        document = Document(
            content=full_markdown,
            metadata=metadata.copy(),
            document_type = 'ppt',
            source_path = file_path,
            title = Path(file_path).name,
            url = f"ppt://{Path(file_path).name}"
        )
        document.metadata['total_slides'] = total_slides
        document.metadata['word_count'] = len(full_markdown)
        document.metadata['processing_method'] = 'python_pptx_vision_combined'
        document.metadata['format'] = 'markdown'
        document.metadata['chunk_strategy'] = 'page'

        logger.info(f"视觉模型处理完成，生成完整Markdown文档，共 {len(full_markdown)} 字符")
        return [document]

    async def _extract_text_only(self, slides_content: List[dict], total_slides: int) -> str:
        """提取纯文本内容（无视觉模型时使用）"""
        logger.info("使用纯文本提取模式")

        markdown_parts = []
        for i, slide_info in enumerate(slides_content):
            markdown_parts.append(f"# 第{i+1}页\n")

            # 添加文本
            for text in slide_info['texts']:
                markdown_parts.append(f"{text}\n")

            # 添加表格
            for table in slide_info['tables']:
                if table:
                    markdown_parts.append("\n")
                    # 表头
                    headers = " | ".join(table[0])
                    markdown_parts.append(f"| {headers} |\n")
                    # 分隔线
                    separators = " | ".join(["---"] * len(table[0]))
                    markdown_parts.append(f"| {separators} |\n")
                    # 数据行
                    for row in table[1:]:
                        row_text = " | ".join(row)
                        markdown_parts.append(f"| {row_text} |\n")

            markdown_parts.append("\n")

        return "".join(markdown_parts)

