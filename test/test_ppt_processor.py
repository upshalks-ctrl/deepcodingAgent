#!/usr/bin/env python3
"""
PPT处理器测试
测试PPTProcessor的功能，包括：
- v1模式: python-pptx后端
- v2模式: 视觉模型 + 内存buffer + base64编码
"""

import asyncio
import sys
import os
from pathlib import Path

# 添加项目根目录到路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from src.document_processors.ppt_processor import PPTProcessor
from src.models.document import Document


class TestPPTProcessor:
    """PPT处理器测试类"""

    def __init__(self):
        self.test_dir = Path(__file__).parent / "test_data"
        self.test_dir.mkdir(exist_ok=True)
        self.test_ppt = self.test_dir / "test_presentation.pptx"

    async def create_test_ppt(self):
        """创建测试PPT文件"""
        if self.test_ppt.exists():
            return

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # 幻灯片1：标题页
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            title1 = slide1.shapes.title
            subtitle1 = slide1.placeholders[1]
            title1.text = "PPT Processor Test"
            subtitle1.text = "Testing PowerPoint Document Processing"

            # 幻灯片2：内容页
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            title2 = slide2.shapes.title
            title2.text = "Overview"
            content2 = slide2.placeholders[1]
            content2.text = """• Introduction to PPT processing
• System architecture
• Implementation details
• Performance considerations"""

            # 幻灯片3：列表页
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            title3 = slide3.shapes.title
            title3.text = "Key Features"
            content3 = slide3.placeholders[1]
            content3.text = """1. Text extraction from slides
2. Shape and text box parsing
3. Table and chart detection
4. Multi-language support
5. Batch processing capability"""

            # 幻灯片4：章节页
            slide4 = prs.slides.add_slide(prs.slide_layouts[1])
            title4 = slide4.shapes.title
            title4.text = "Technical Details"
            content4 = slide4.placeholders[1]
            content4.text = """Chapter 1: Data Structure
- Slide model definition
- Text box hierarchy
- Shape relationships

Chapter 2: Processing Pipeline
- Input validation
- Content extraction
- Output formatting"""

            prs.save(str(self.test_ppt))
            print(f"✓ Created test PPT: {self.test_ppt}")
        except Exception as e:
            print(f"✗ Failed to create test PPT: {e}")
            raise

    async def test_v1_processor(self):
        """测试PPTProcessor v1 - python-pptx方法"""
        print("\n" + "="*70)
        print("TEST 1: PPTProcessor v1 - python-pptx method")
        print("="*70)

        processor = PPTProcessor()

        try:
            documents = await processor.process(str(self.test_ppt))
            print(f"✓ python-pptx processed {len(documents)} slides")

            for i, doc in enumerate(documents, 1):
                print(f"\n  Slide {i}:")
                print(f"    Title: {doc.title}")
                print(f"    Slide number: {doc.metadata.get('slide_number', 'N/A')}")
                print(f"    Slide title: {doc.metadata.get('slide_title', 'N/A')}")
                print(f"    Content length: {len(doc.content)} chars")
                print(f"    Preview: {doc.content[:100]}...")

            return True
        except Exception as e:
            print(f"✗ python-pptx method failed: {e}")
            import traceback
            traceback.print_exc()
            return False

    async def test_v1_metadata(self):
        """测试PPTProcessor v1 - 元数据"""
        print("\n" + "="*70)
        print("TEST 2: PPTProcessor v1 - Metadata")
        print("="*70)

        processor = PPTProcessor()

        try:
            documents = await processor.process(str(self.test_ppt))

            if not documents:
                print("✗ No documents returned")
                return False

            doc = documents[0]
            required_fields = ['source_path', 'title', 'url']

            print(f"✓ Document has {len(doc.metadata)} metadata fields")

            for field in required_fields:
                if field in doc.metadata:
                    print(f"  ✓ {field}: {doc.metadata[field]}")
                else:
                    print(f"  ✗ Missing {field}")

            # 检查PPT特定字段
            ppt_fields = ['slide_number', 'slide_title', 'total_slides', 'word_count']
            for field in ppt_fields:
                if field in doc.metadata:
                    print(f"  ✓ {field}: {doc.metadata[field]}")

            return True
        except Exception as e:
            print(f"✗ Metadata test failed: {e}")
            return False

    async def test_v2_processor_basic(self):
        """测试PPTProcessor v2模式 - 基础功能"""
        print("\n" + "="*70)
        print("TEST 3: PPTProcessor v2 - Basic Processing")
        print("="*70)

        processor = PPTProcessor()

        try:
            # 使用v2模式，不使用视觉模型（会回退到python-pptx）
            documents = await processor.process(
                str(self.test_ppt),
                mode='vision_memory_buffer',
                chunk_after=False
            )
            print(f"✓ v2模式处理了 {len(documents)} 页")

            for i, doc in enumerate(documents, 1):
                print(f"\n  Document {i}:")
                print(f"    Title: {doc.title}")
                print(f"    Content length: {len(doc.content)} chars")
                print(f"    Total slides: {doc.metadata.get('total_slides', 'N/A')}")
                print(f"    Processing method: {doc.metadata.get('processing_method', 'N/A')}")

            return True
        except Exception as e:
            print(f"✗ Basic processing failed: {e}")
            import traceback
            traceback.print_exc()
            return False

    async def test_v2_processor_with_vision(self):
        """测试PPTProcessor v2模式 - 视觉模型功能"""
        print("\n" + "="*70)
        print("TEST 4: PPTProcessor v2 - With Vision Model")
        print("="*70)

        processor = PPTProcessor()

        try:
            documents = await processor.process(
                str(self.test_ppt),
                mode='vision_memory_buffer',
                dpi=150,
                chunk_after=True,
                chunk_size=500,
                chunk_overlap=50
            )
            print(f"✓ 视觉模型处理创建了 {len(documents)} 个块")

            for i, doc in enumerate(documents, 1):
                print(f"\n  Chunk {i}:")
                print(f"    Title: {doc.title}")
                print(f"    Content length: {len(doc.content)} chars")
                print(f"    Chunk type: {doc.metadata.get('chunk_type', 'N/A')}")
                print(f"    Chunk title: {doc.metadata.get('chunk_title', 'N/A')}")
                print(f"    Processing method: {doc.metadata.get('processing_method', 'N/A')}")

            return True
        except Exception as e:
            print(f"✗ Vision processing failed: {e}")
            print(f"  Note: PPT v2 requires PowerPoint application (Windows)")
            print(f"  This is expected if PowerPoint is not available")
            import traceback
            traceback.print_exc()
            return False

    async def test_v2_memory_buffer(self):
        """测试PPTProcessor v2模式 - 内存buffer功能"""
        print("\n" + "="*70)
        print("TEST 5: PPTProcessor v2 - Memory Buffer & Base64")
        print("="*70)

        processor = PPTProcessor()

        try:
            # 测试内部方法
            documents = await processor._process_with_vision_memory(
                str(self.test_ppt),
                dpi=150,
                chunk_after=False
            )
            print(f"✓ Memory buffer processing created {len(documents)} documents")

            doc = documents[0]
            print(f"\n  Complete Document:")
            print(f"    Title: {doc.title}")
            print(f"    Content length: {len(doc.content)} chars")
            print(f"    Total slides: {doc.metadata.get('total_slides', 'N/A')}")
            print(f"    Processing method: {doc.metadata.get('processing_method', 'N/A')}")
            print(f"    Format: {doc.metadata.get('format', 'N/A')}")
            print(f"    Preview: {doc.content[:200]}...")

            return True
        except Exception as e:
            print(f"✗ Memory buffer processing failed: {e}")
            print(f"  Note: PPT v2 requires PowerPoint application (Windows)")
            print(f"  This is expected if PowerPoint is not available")
            import traceback
            traceback.print_exc()
            return False

    async def test_text_extraction_quality(self):
        """测试文本提取质量"""
        print("\n" + "="*70)
        print("TEST 6: Text Extraction Quality")
        print("="*70)

        processor = PPTProcessor()

        try:
            documents = await processor.process(str(self.test_ppt))

            # 检查是否提取了预期内容
            all_content = " ".join(doc.content for doc in documents)

            expected_keywords = [
                "PPT Processor Test",
                "Overview",
                "Key Features",
                "Technical Details"
            ]

            found_keywords = []
            for keyword in expected_keywords:
                if keyword in all_content:
                    found_keywords.append(keyword)
                    print(f"  ✓ Found keyword: {keyword}")
                else:
                    print(f"  ✗ Missing keyword: {keyword}")

            print(f"\n  Found {len(found_keywords)}/{len(expected_keywords)} keywords")

            # 检查文本完整性
            total_length = sum(len(doc.content) for doc in documents)
            print(f"  Total content length: {total_length} chars")

            return len(found_keywords) >= len(expected_keywords) * 0.75
        except Exception as e:
            print(f"✗ Text extraction test failed: {e}")
            return False

    async def test_slide_segmentation(self):
        """测试幻灯片分割"""
        print("\n" + "="*70)
        print("TEST 7: Slide Segmentation")
        print("="*70)

        processor = PPTProcessor()

        try:
            documents = await processor.process(str(self.test_ppt))

            # 检查每个幻灯片是否被正确分割
            expected_slides = 4
            print(f"  Expected slides: {expected_slides}")
            print(f"  Actual documents: {len(documents)}")

            if len(documents) >= expected_slides:
                print(f"  ✓ Correct number of slides")

                # 检查幻灯片编号
                slide_numbers = [doc.metadata.get('slide_number') for doc in documents]
                print(f"  Slide numbers: {slide_numbers}")

                # 检查每页内容的非空性
                non_empty = sum(1 for doc in documents if doc.content.strip())
                print(f"  Non-empty slides: {non_empty}/{len(documents)}")

                return non_empty == len(documents)
            else:
                print(f"  ✗ Incorrect number of slides")
                return False

        except Exception as e:
            print(f"✗ Slide segmentation test failed: {e}")
            return False

    async def run_all_tests(self):
        """运行所有测试"""
        print("\n" + "="*70)
        print("PPT PROCESSOR TEST SUITE")
        print("="*70)

        # 创建测试PPT
        await self.create_test_ppt()

        # 运行测试
        results = []
        results.append(("v1 python-pptx", await self.test_v1_processor()))
        results.append(("v1 Metadata", await self.test_v1_metadata()))
        results.append(("v2 Basic", await self.test_v2_processor_basic()))
        results.append(("v2 Vision", await self.test_v2_processor_with_vision()))
        results.append(("v2 Memory Buffer", await self.test_v2_memory_buffer()))
        results.append(("Text Quality", await self.test_text_extraction_quality()))
        results.append(("Slide Segmentation", await self.test_slide_segmentation()))

        # 总结
        print("\n" + "="*70)
        print("TEST SUMMARY")
        print("="*70)

        passed = sum(1 for _, result in results if result)
        total = len(results)

        for test_name, result in results:
            status = "✓ PASS" if result else "✗ FAIL"
            print(f"  {status} - {test_name}")

        print(f"\n  Total: {passed}/{total} tests passed")

        return passed == total


async def main():
    """主函数"""
    test_suite = TestPPTProcessor()
    success = await test_suite.run_all_tests()

    if success:
        print("\n🎉 All tests passed!")
    else:
        print("\n⚠️  Some tests failed")

    print("\n" + "="*70)
    print("NOTES")
    print("="*70)
    print("• PPTProcessor v2 mode uses PowerPoint application (Windows)")
    print("• Vision model functionality needs vision LLM to be configured")
    print("• Memory buffer tests require COM interfaces to be available")
    print("• Supports mode='vision_memory_buffer' for advanced processing")
    print("• Supports automatic page-based chunking")

    return success


if __name__ == "__main__":
    import logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    result = asyncio.run(main())
    sys.exit(0 if result else 1)
